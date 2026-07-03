#!/usr/bin/env python3
"""Generate PowerPoint presentation for Sprintify diploma thesis."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE_DARK = RGBColor(58, 4, 255)  # #3A04FF - Logo color
BLUE_LIGHT = RGBColor(99, 102, 241)  # Indigo
WHITE = RGBColor(255, 255, 255)
GRAY_DARK = RGBColor(30, 41, 59)  # Slate 800
GRAY_LIGHT = RGBColor(241, 245, 249)  # Slate 100


def add_title_slide(title, subtitle=None):
    """Add a title slide with dark background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = GRAY_DARK
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(148, 163, 184)  # Slate 400
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)

    return slide


def add_content_slide(title, bullets, note=None):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # White background (default)

    # Accent bar at top
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE_DARK
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GRAY_DARK

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullets (lines starting with spaces or tabs)
        if bullet.startswith("  ") or bullet.startswith("\t"):
            p.text = "  " + bullet.strip()
            p.level = 1
            p.font.size = Pt(18)
        else:
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)

        p.font.color.rgb = RGBColor(51, 65, 85)  # Slate 700
        p.space_before = Pt(8)
        p.space_after = Pt(4)

    # Add speaker note if provided
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
    """Add a slide with two columns."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE_DARK
    accent.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GRAY_DARK

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(5.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.9), Inches(5.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(6)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(6)

    return slide


def add_section_slide(title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BLUE_DARK
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


# ============================================================================
# CREATE SLIDES
# ============================================================================

# Slide 1: Title
add_title_slide(
    "Sprintify",
    "Entwicklung eines Capacity-Planning-Tools mit Jira-Integration\n\nTimo Haldi | HF ICT 2023 F | Coach: Ardin Ibraimi"
)

# Slide 2: Agenda
add_content_slide("Agenda", [
    "1. Vorstellung und Kontext",
    "2. Problemstellung und Auftrag",
    "3. Vorgehen und Projektplanung",
    "4. Architektur und Technologieentscheidungen",
    "5. Ergebnisse: Die Module im Detail",
    "6. Sprints und Entwicklungsprozess",
    "7. Live-Demo",
    "8. Kritische Würdigung",
    "9. Persönliche Reflexion",
    "10. Fazit und Empfehlung",
])

# Slide 3: Vorstellung
add_two_column_slide(
    "Vorstellung und Kontext",
    "Wer bin ich?",
    [
        "Cloud Engineer bei der Netcloud AG",
        "Abteilung Public Cloud Solutions",
        "Tagesgeschäft: Azure-Kundenprojekte",
        "Team arbeitet nach Scrum",
        "Zweiwöchige Sprints",
    ],
    "Was macht die Netcloud AG?",
    [
        "Schweizer IT-Dienstleister",
        "Fokus auf Cloud-Infrastruktur",
        "Kunden aus verschiedenen Branchen",
        "Unterschiedliche Jira-Setups",
    ]
)

# Slide 4: Problemstellung
add_content_slide("Problemstellung", [
    "Sprint-Planung vor Sprintify:",
    "  Kapazitätsplanung via Excel-Tabellen",
    "  Verstreut über SharePoint-Ordner",
    "  Kein zentraler Ort für historische Velocity-Daten",
    "  Abwesenheiten fliessen nur informell in die Planung ein",
    "",
    "Die Frage \"Wie viele Story Points schaffen wir?\" konnte niemand datenbasiert beantworten",
    "",
    "Konsequenzen:",
    "  Regelmässig zu viel oder zu wenig eingeplant",
    "  Sprint-Ziele verfehlt",
    "  Mehr Abstimmungsaufwand, weniger Planungssicherheit",
], "Konkretes Beispiel: In einem Sprint haben wir 40 SP geplant, zwei Leute waren aber im Urlaub. Am Ende haben wir 22 SP geschafft.")

# Slide 5: Auftrag
add_content_slide("Auftrag", [
    "Was sollte Sprintify leisten?",
    "  1. Sprints, User Stories und Teammitglieder aus Jira übernehmen",
    "  2. Kapazitätsplanung pro Person und Sprint (wochenweise)",
    "  3. Sprint Analytics: Burndown, Velocity, Scope-Changes",
    "  4. SP-Empfehlung basierend auf Kapazität und Velocity",
    "  5. Rollenbasierte Zugriffskontrolle",
    "",
    "Rahmenbedingungen:",
    "  Deployment auf Azure App Service",
    "  Ein-Mann-Projekt",
    "  12 Monate Projektlaufzeit",
])

# Slide 6: Marktanalyse
add_content_slide("Marktanalyse", [
    "Tempo Timesheets",
    "  + Zeiterfassung, Jira-Plugin",
    "  - Keine wochenbasierte Kapazitätsplanung",
    "",
    "Jira Advanced Roadmaps",
    "  + Portfolio-Planung",
    "  - Teuer, komplex, keine Team-Kapazitätsansicht",
    "",
    "Forecast.app",
    "  + Projektplanung",
    "  - Keine direkte Jira-Sync, keine SP-Empfehlung",
    "",
    "Keines verbindet Kapazitätsplanung + Jira-Sync + Sprint Analytics",
])

# Slide 7: Vorgehen
add_content_slide("Vorgehen", [
    "Methodik: Iteratives Vorgehen (angelehnt an Scrum)",
    "  Zweiwöchige Iterationen mit definierten Zielen",
    "  Regelmässige Besprechung mit dem Auftraggeber",
    "",
    "Vier Projektphasen:",
    "  Phase 1 (Monat 1-2): Analyse & Konzeption",
    "  Phase 2 (Monat 3-5): Backend-Entwicklung",
    "  Phase 3 (Monat 5-8): Frontend-Entwicklung",
    "  Phase 4 (Monat 9-11): Integration & Deployment",
], "Erklären, warum Scrum als Einzelperson trotzdem Sinn ergibt: feste Iterationen erzwingen regelmässige Standortbestimmung.")

# Section: Architektur
add_section_slide("Architektur")

# Slide 8: Architekturübersicht
add_content_slide("Architekturübersicht", [
    "Frontend:",
    "  React 19, TypeScript, Vite, Tailwind CSS",
    "",
    "Backend:",
    "  Node.js, Express.js, Sequelize ORM",
    "  10 Route-Module, Middleware-Stack",
    "",
    "Datenbank:",
    "  PostgreSQL (Managed auf Azure)",
    "",
    "Externe Systeme:",
    "  Jira Cloud (REST API v2, Agile API v1.0)",
    "  Microsoft Entra ID (MSAL)",
])

# Slide 9: Datenmodell
add_content_slide("Datenmodell", [
    "7 Modelle:",
    "",
    "  Project → Sprint → UserStory",
    "  Project → ProjectUser ↔ User",
    "  Sprint → CapacityPlan ↔ User",
    "  Sprint → Retrospective",
    "",
    "Wichtige Beziehungen:",
    "  ProjectUser: Many-to-Many mit Rollen (admin, member, viewer)",
    "  CapacityPlan: Pro User und Sprint, wochenweise Kapazitätsdaten",
])

# Slide 10: Sicherheit
add_two_column_slide(
    "Sicherheit und Zugriffskontrolle",
    "Authentifizierung",
    [
        "Microsoft Entra ID mit MSAL",
        "OAuth 2.0 Authorization Code Flow",
        "Fallback: lokale JWT-Authentifizierung",
        "Tokens in HTTP-only Cookies",
    ],
    "Autorisierung",
    [
        "Global: Admin, Member",
        "Pro Projekt: Admin, Member, Viewer",
        "",
        "Sicherheits-Middleware:",
        "• CSRF-Schutz",
        "• Rate Limiting",
        "• Input-Sanitisierung",
    ]
)

# Section: Jira
add_section_slide("Jira-Integration")

# Slide 11: Jira-Integration
add_content_slide("Jira-Integration", [
    "Die grösste technische Herausforderung",
    "",
    "Zwei APIs:",
    "  REST API v2: Issues (User Stories)",
    "  Agile API v1.0: Boards und Sprints",
    "",
    "Synchronisation in drei Schritten:",
    "  1. Boards und Sprints abrufen",
    "  2. Issues pro Sprint synchronisieren",
    "  3. Jira-Benutzer mit lokalen Benutzern verknüpfen",
    "",
    "Besonderheiten:",
    "  Konfigurierbares Story-Point-Feld",
    "  Status-Mapping auf einheitliches Schema",
    "  Automatische Sync alle 15 Minuten",
], "Der schwierigste Teil war nicht die API, sondern die unterschiedlichen Konfigurationen bei verschiedenen Kunden.")

# Section: Module
add_section_slide("Module im Detail")

# Slide 12: Capacity Planning
add_content_slide("Modul: Capacity Planning", [
    "Funktionen:",
    "  Wochenweise Aufschlüsselung pro Teammitglied",
    "  4 Kategorien: Holiday, Customer, Internal, Other",
    "  Aggregierte Ansicht: Gesamtkapazität, Durchschnitt",
    "  Inline-Editing direkt in der Tabelle",
    "",
    "SP-Empfehlung (neu):",
    "  Formel: empfohlene SP = Avg Velocity × Kapazitätsfaktor",
    "  Kapazitätsfaktor = aktuelle / historische Durchschnittskapazität",
    "  Berücksichtigt die letzten 6 abgeschlossenen Sprints",
], "Wenn zwei Leute im Urlaub sind, sinkt die Kapazität. Die Empfehlung passt die SP automatisch nach unten an.")

# Slide 13: Sprint Analytics
add_content_slide("Modul: Sprint Analytics", [
    "Vier Ansichten:",
    "",
    "  Overview: Kennzahlen + Burndown-Chart",
    "  → Sind wir auf Kurs?",
    "",
    "  Compare: Velocity vs. vorheriger Sprint",
    "  → Werden wir besser oder schlechter?",
    "",
    "  Changes: Scope-Änderungen während des Sprints",
    "  → Warum haben wir das Ziel verfehlt?",
    "",
    "  Team Performance: Velocity pro Person (Balkendiagramm)",
    "  → Ist die Arbeit gleichmässig verteilt?",
])

# Slide 14: Sprint History & Forecast
add_two_column_slide(
    "Modul: Sprint History & Forecast",
    "Sprint History",
    [
        "Velocity über alle Sprints",
        "Darstellung als Balkendiagramm",
        "Grundlage für Schätzungen",
    ],
    "Forecast",
    [
        "Verbleibende Backlog-Punkte",
        "Geschätzte Anzahl Sprints",
        "  • Best Case",
        "  • Average Case",
        "  • Worst Case",
        "Geschätztes Fertigstellungsdatum",
    ]
)

# Section: Entwicklung
add_section_slide("Entwicklungsprozess")

# Slide 15: Sprints
add_content_slide("Sprints und Phasen", [
    "Phase 1 - Analyse (Sprint 1-4):",
    "  Anforderungen, Marktanalyse, Architekturentscheid",
    "",
    "Phase 2 - Backend (Sprint 5-10):",
    "  Datenmodell, API, Jira-Integration",
    "",
    "Phase 3 - Frontend (Sprint 11-18):",
    "  Dashboard, Kapazitätsplanung, Analytics",
    "",
    "Phase 4 - Deployment (Sprint 19-24):",
    "  Azure, Tests, Performance, SP-Empfehlung",
], "Die Jira-Integration hat mehr Sprints gebraucht als geplant. Die SP-Empfehlung kam als Anforderung erst spät dazu.")

# Slide 16: Deployment
add_content_slide("Deployment und Betrieb", [
    "Azure App Service:",
    "  Kein eigener Server",
    "  Zentrale Updates",
    "  Automatische Skalierung",
    "",
    "Komponenten:",
    "  Backend als App Service",
    "  Frontend als pre-built Static Files",
    "  PostgreSQL als Managed Database",
    "",
    "Konfiguration:",
    "  Umgebungsvariablen über Azure App Settings",
    "  Datenbankmigrationen via Sequelize CLI",
])

# Slide 17: Live-Demo
add_section_slide("Live-Demo")

add_content_slide("Demo-Ablauf", [
    "1. Login (Microsoft Entra ID)",
    "",
    "2. Dashboard - Projektübersicht, Kennzahlen",
    "",
    "3. Kapazitätsplanung",
    "   Wochenweise Eingabe, Kategorien, SP-Empfehlung",
    "",
    "4. Sprint Analytics",
    "   Burndown-Chart, Velocity-Vergleich",
    "",
    "5. Scope Changes",
    "",
    "6. Sprint History - Velocity-Balkendiagramm",
], "Demo vorbereiten und vorher durchspielen. Screenshots als Backup falls Jira/Azure nicht erreichbar.")

# Slide 18: Kritische Würdigung
add_two_column_slide(
    "Kritische Würdigung",
    "Was gut gelaufen ist",
    [
        "Iteratives Vorgehen hat sich bewährt",
        "Anforderungen haben sich geändert",
        "App läuft mit Produktivdaten",
        "SP-Empfehlung ist einzigartig",
    ],
    "Was ich anders machen würde",
    [
        "Jira-Integration: Spike zu Beginn",
        "Vielfalt der Konfigurationen unterschätzt",
        "Automatisierte Tests nur punktuell",
        "Single-Tenant Architektur limitiert",
    ]
)

# Slide 19: Persönliche Reflexion
add_content_slide("Persönliche Reflexion", [
    "Was ich gelernt habe:",
    "  Die grösste Herausforderung steckt in externen Systemen",
    "  Allein arbeiten: schnelle Entscheidungen, aber keine Reviews",
    "  Ohne Iterationen hätte ich die Änderungen nicht geschafft",
    "  Vorher Frontend/Infra - jetzt auch Backend-Erfahrung",
    "",
    "Was ich mitnehme:",
    "  Integrationen immer mit einem Spike starten",
    "  Automatisierte Tests von Anfang an",
    "  Auch als Solo-Entwickler Code-Reviews simulieren",
])

# Slide 20: Empfehlung
add_content_slide("Empfehlung und nächste Schritte", [
    "Empfehlung: Sprintify produktiv einsetzen",
    "",
    "Nächste Schritte:",
    "  1. Jira-Anbindung über Service-Account",
    "  2. Einführungsschulung für das Team",
    "  3. Feedback sammeln in den ersten drei Sprints",
    "  4. Automatisierte Tests ausbauen",
    "  5. Mittelfristig: Mobile Ansicht für Stand-ups",
    "",
    "Ohne Sprintify:",
    "  Sprint-Planung bleibt auf Excel angewiesen",
    "  Schätzungen basieren auf Bauchgefühl",
])

# Slide 21: Abschluss
add_title_slide(
    "Sprintify",
    "Gibt dem Team eine datenbasierte Antwort auf die Frage:\n\"Wie viele Story Points schaffen wir im nächsten Sprint?\"\n\nFragen?"
)

# Save presentation
output_path = "/Users/timohaldi/Documents/Github/Private/timohone/sprintify/docs/Sprintify_Praesentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
