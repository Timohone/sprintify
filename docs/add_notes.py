#!/usr/bin/env python3
"""Add speaker notes to every slide of the Sprintify presentation."""

from pptx import Presentation

prs = Presentation("/Users/timohaldi/Documents/Github/Private/timohone/sprintify/docs/Sprintify_Praesentation.pptx")

notes = {
    0: """Guten Tag, mein Name ist Timo Haldi. Ich präsentiere heute meine Diplomarbeit: Sprintify – ein Capacity-Planning-Tool mit Jira-Integration.
Ich studiere an der SIW Höheren Fachschule für Wirtschaft und Informatik im Lehrgang HF ICT 2023 F. Mein Coach ist Ardin Ibraimi.""",

    1: """Ich habe die Präsentation in zehn Blöcke aufgeteilt. Wir starten mit meiner Vorstellung und dem Kontext, gehen dann auf die Problemstellung ein, schauen uns die Architektur und die einzelnen Module an, und am Ende gibt es eine Live-Demo. Danach folgen die kritische Würdigung und meine persönliche Reflexion.""",

    2: """Ich arbeite als Cloud Engineer bei der Netcloud AG in der Abteilung Public Cloud Solutions. Im Tagesgeschäft betreue ich Azure-Kundenprojekte, Infrastruktur und Managed Services. Unser Team arbeitet nach Scrum mit zweiwöchigen Sprints.

Die Netcloud AG ist ein Schweizer IT-Dienstleister mit Fokus auf Cloud-Infrastruktur. Wir haben Kunden aus verschiedenen Branchen, die jeweils ihre eigenen Jira-Setups mitbringen. Das ist später noch relevant, wenn es um die Jira-Integration geht.""",

    3: """Bevor es Sprintify gab, haben wir die Sprint-Planung mit Excel-Tabellen gemacht. Diese Tabellen lagen verstreut in SharePoint-Ordnern. Es gab keinen zentralen Ort für historische Velocity-Daten. Abwesenheiten wie Ferien, Teilzeit oder Kundenprojekte sind nur informell in die Planung eingeflossen.

Die zentrale Frage – wie viele Story Points schaffen wir im nächsten Sprint – konnte niemand datenbasiert beantworten.

Ein konkretes Beispiel: In einem Sprint haben wir 40 Story Points eingeplant. Zwei Leute waren aber im Urlaub. Am Ende haben wir nur 22 Story Points geschafft. Das ist genau das Problem, das Sprintify lösen soll.""",

    4: """Was sollte Sprintify konkret leisten? Erstens: Sprints, User Stories und Teammitglieder automatisch aus Jira übernehmen. Zweitens: eine Kapazitätsplanung pro Person und Sprint, aufgeteilt nach Wochen und Kategorien. Drittens: Sprint Analytics mit Burndown-Charts, Velocity-Verlauf und Scope-Change-Tracking. Viertens: eine Story-Point-Empfehlung basierend auf Kapazität und historischer Velocity. Und fünftens: rollenbasierte Zugriffskontrolle.

Die Rahmenbedingungen: Deployment auf Azure App Service, ein Ein-Mann-Projekt über 12 Monate.""",

    5: """Ich habe bestehende Lösungen evaluiert. Tempo Timesheets bietet Zeiterfassung als Jira-Plugin, hat aber keine wochenbasierte Kapazitätsplanung. Jira Advanced Roadmaps ist gut für Portfolio-Planung, aber teuer, komplex und bietet keine Team-Kapazitätsansicht. Forecast.app hat Projektplanung, aber keine direkte Jira-Synchronisation und keine SP-Empfehlung.

Keines dieser Tools verbindet wochenbasierte Kapazitätsplanung mit automatischer Jira-Sync und Sprint Analytics in einer Oberfläche. Für kleinere Teams wie unseres sind sie ausserdem oft zu teuer.""",

    6: """Für das Vorgehen habe ich mich für ein iteratives Vorgehen angelehnt an Scrum entschieden. Zweiwöchige Iterationen mit definierten Zielen und regelmässige Besprechung mit dem Auftraggeber.

Warum Scrum als Einzelperson? Feste Iterationen erzwingen eine regelmässige Standortbestimmung. Ich musste alle zwei Wochen prüfen: Bin ich noch auf Kurs? Stimmen die Prioritäten noch?

Das Projekt hatte vier Phasen: Analyse und Konzeption in den ersten zwei Monaten, Backend-Entwicklung von Monat drei bis fünf, Frontend-Entwicklung von Monat fünf bis acht, und Integration und Deployment von Monat neun bis elf.

Die Anforderungen haben sich im Projektverlauf geändert, vor allem bei der Jira-Integration. Genau dafür war das iterative Vorgehen richtig.""",

    7: """Wir kommen jetzt zum Architektur-Block. Hier schauen wir uns den Technologie-Stack, das Datenmodell und die Sicherheitskonzepte an.""",

    8: """Die Architektur ist eine klassische Client-Server-Architektur. Das Frontend ist eine React 19 Single Page Application mit TypeScript, gebaut mit Vite und gestylt mit Tailwind CSS.

Das Backend läuft auf Node.js mit Express.js und Sequelize als ORM. Es gibt 10 Route-Module und einen Middleware-Stack mit JWT-Authentifizierung, CSRF-Schutz und Rate Limiting.

Als Datenbank verwende ich PostgreSQL. Extern wird Jira Cloud angebunden über zwei APIs: die REST API v2 für Issues und die Agile API v1.0 für Boards und Sprints. Die Authentifizierung läuft über Microsoft Entra ID mit MSAL.

Warum diese Technologien? Node.js wegen meiner vorhandenen JavaScript-Erfahrung und weil asynchrone Verarbeitung gut zur Jira-API passt. PostgreSQL wegen der relationalen Datenstruktur – wir haben Many-to-Many-Beziehungen zwischen Usern und Projekten. React wegen der komponentenbasierten Architektur und der grossen Community.""",

    9: """Das Datenmodell besteht aus sieben Modellen. Ein Project hat mehrere Sprints, und ein Sprint hat mehrere User Stories. Projekte und User sind über die ProjectUser-Tabelle verbunden – das ist eine Many-to-Many-Beziehung mit Rollen: Admin, Member und Viewer.

Pro Sprint und User gibt es einen CapacityPlan mit wochenweisen Kapazitätsdaten. Und Sprints können Retrospectives haben.

Dieses Datenmodell ist die Grundlage für alles, was jetzt kommt – die Module bauen darauf auf.""",

    10: """Bei der Sicherheit gibt es zwei Ebenen. Die Authentifizierung läuft über Microsoft Entra ID mit MSAL und dem OAuth 2.0 Authorization Code Flow. Als Fallback gibt es eine lokale JWT-basierte Authentifizierung. Tokens werden in HTTP-only Cookies gespeichert.

Die Autorisierung funktioniert auf zwei Ebenen: Global gibt es Admin und Member. Pro Projekt gibt es Admin, Member und Viewer – das wird über die ProjectUser-Tabelle gesteuert.

Dazu kommt eine Sicherheits-Middleware mit CSRF-Schutz, Rate Limiting und Input-Sanitisierung.

Die App ist Single-Tenant: eine Instanz pro Organisation. Für den aktuellen Einsatz bei der Netcloud AG brauchen wir keine Mandantentrennung.""",

    11: """Wir kommen jetzt zur Jira-Integration – das war die grösste technische Herausforderung des Projekts.""",

    12: """Die Jira-Integration nutzt zwei separate APIs. Die REST API v2 für Issues, also User Stories, und die Agile API v1.0 für Boards und Sprints.

Die Synchronisation läuft in drei Schritten: Zuerst Boards und Sprints abrufen, dann Issues pro Sprint synchronisieren, und zuletzt die Jira-Benutzer mit lokalen Benutzern verknüpfen.

Ein paar Besonderheiten: Das Story-Point-Feld ist konfigurierbar, weil verschiedene Kunden verschiedene Custom Fields verwenden. Der Standard ist customfield_10016. Es gibt ein Status-Mapping, das die Jira-Status auf ein einheitliches Schema abbildet: To Do, In Progress und Done. Jedes Projekt kann seine eigene Jira-Instanz haben. Und der Scheduler synchronisiert automatisch alle 15 Minuten.

Ein Fallstrick war die Synchronisationsreihenfolge: Stories können in mehreren Sprints vorkommen. Man muss chronologisch aufsteigend synchronisieren, damit die Zuordnung zum neuesten Sprint gewinnt.

Das war der schwierigste Teil des Projekts. Nicht wegen der API selbst, sondern wegen der unterschiedlichen Konfigurationen bei verschiedenen Kunden. Zum Beispiel hatte ein Kunde die Story Points in einem anderen Custom Field als der Standard.""",

    13: """Jetzt kommen die Module im Detail. Hier zeige ich, was Sprintify konkret kann.""",

    14: """Das Kapazitätsplanungs-Modul ist das Herzstück von Sprintify. Es bietet eine wochenweise Aufschlüsselung pro Teammitglied in vier Kategorien: Holiday, Customer, Internal und Other. Team-Leads können die Pläne direkt in der Tabelle bearbeiten. Es gibt eine aggregierte Ansicht mit der Gesamtkapazität und dem Durchschnitt pro Person.

Neu dazugekommen ist die SP-Empfehlung. Die Formel lautet: Empfohlene Story Points gleich durchschnittliche Velocity mal Kapazitätsfaktor. Der Kapazitätsfaktor ist die aktuelle Sprintkapazität geteilt durch die historische Durchschnittskapazität der letzten sechs abgeschlossenen Sprints.

Um das greifbar zu machen: Wenn normalerweise acht Leute im Team sind und im nächsten Sprint zwei im Urlaub, sinkt die Kapazität um 25 Prozent. Die Empfehlung passt die Story Points automatisch nach unten an. Das war vorher Bauchgefühl.""",

    15: """Sprint Analytics hat vier Ansichten. Die Overview-Ansicht zeigt die Kennzahlen des aktuellen Sprints und ein Burndown-Chart – ideal versus tatsächlich. Die Frage, die sie beantwortet: Sind wir auf Kurs?

Die Compare-Ansicht vergleicht die aktuelle Velocity mit dem vorherigen Sprint und zeigt Langzeit-Benchmarks. Die Frage: Werden wir besser oder schlechter?

Die Changes-Ansicht zeigt Scope-Änderungen – also Stories, die während des Sprints hinzugekommen sind. Die Frage: Warum haben wir das Sprint-Ziel verfehlt?

Und die Team-Performance-Ansicht zeigt Velocity und Completion Rate pro Person als Balkendiagramm. Die Frage: Ist die Arbeit gleichmässig verteilt?""",

    16: """Die Sprint History zeigt die Velocity über alle abgeschlossenen Sprints als Balkendiagramm. Das ist die Grundlage für datenbasierte Schätzungen.

Der Forecast berechnet daraus, wie viele Sprints noch nötig sind, um das Backlog abzuarbeiten. Er zeigt drei Szenarien: Best Case, Average Case und Worst Case. Dazu gibt es ein geschätztes Fertigstellungsdatum.""",

    17: """Wir schauen uns jetzt den Entwicklungsprozess an – wie das Projekt über die 12 Monate verlaufen ist.""",

    18: """Das Projekt lief über 24 zweiwöchige Sprints. In Phase eins, Sprint eins bis vier, ging es um Anforderungen, Marktanalyse und den Architekturentscheid.

In Phase zwei, Sprint fünf bis zehn, habe ich das Backend gebaut: Datenmodell, API und die Jira-Integration. Die Jira-Integration hat mehr Sprints gebraucht als geplant – das war der erste grössere Planungsabweichung.

In Phase drei, Sprint elf bis achtzehn, kam das Frontend dran: Dashboard, Kapazitätsplanung und Analytics.

In Phase vier, Sprint neunzehn bis vierundzwanzig, ging es um Azure Deployment, Tests mit Produktivdaten und die SP-Empfehlung. Die SP-Empfehlung kam als Anforderung erst spät dazu. Das zeigt, warum iteratives Vorgehen richtig war: Neue Anforderungen liessen sich in den nächsten Sprint einplanen.""",

    19: """Sprintify läuft auf Azure App Service. Das heisst: kein eigener Server, zentrale Updates und automatische Skalierung.

Das Backend läuft als App Service, das Frontend wird als pre-built Static Files ausgeliefert, und PostgreSQL läuft als Managed Database.

Die Konfiguration erfolgt über Umgebungsvariablen in den Azure App Settings. Im Frontend gibt es ein inject-config.js Script, das die Konfiguration zur Laufzeit injiziert. Datenbankmigrationen laufen über die Sequelize CLI.""",

    20: """Und jetzt zeige ich die Applikation live.""",

    21: """Ich zeige jetzt sechs Bereiche. Zuerst den Login über Microsoft Entra ID. Dann das Dashboard mit Projektübersicht und Kennzahlen. Dann die Kapazitätsplanung mit der wochenweisen Eingabe und der SP-Empfehlung. Danach Sprint Analytics mit dem Burndown-Chart und Velocity-Vergleich. Dann die Scope Changes. Und zum Schluss die Sprint History mit dem Velocity-Balkendiagramm.

Falls Jira oder Azure nicht erreichbar sein sollten, habe ich Screenshots als Backup vorbereitet.""",

    22: """Zur kritischen Würdigung. Was ist gut gelaufen? Das iterative Vorgehen hat sich bewährt. Die Anforderungen haben sich mehrfach geändert, vor allem bei der Jira-Integration, und das konnten wir auffangen. Die Applikation läuft mit Produktivdaten und die Ergebnisse stimmen. Die SP-Empfehlung gab es so in keinem der evaluierten Tools.

Was würde ich anders machen? Bei der Jira-Integration hätte ich zu Beginn einen Spike mit einem Minimal-Prototyp machen sollen. Die Vielfalt der Konfigurationen habe ich unterschätzt. Als Ein-Mann-Projekt fehlten Code-Reviews – automatisierte Tests hätten das teilweise kompensiert, habe ich aber nur punktuell geschrieben. Und die App ist Single-Tenant. Falls sie bei mehreren Kunden laufen soll, müsste man die Mandantentrennung ergänzen.""",

    23: """Zur persönlichen Reflexion. Die grösste Erkenntnis: Bei Integrationsprojekten steckt die Herausforderung selten in der eigenen Codebasis. Sie steckt in den Eigenheiten und Inkonsistenzen des Systems, an das man anbindet.

Allein arbeiten hat zwei Seiten: Schnelle Entscheidungen, aber keine Code-Reviews. Ohne Iterationen hätte ich die wechselnden Anforderungen nicht auffangen können. Vorher hatte ich vor allem Frontend- und Infrastruktur-Erfahrung. Durch das Projekt kann ich jetzt auch Backend.

Für zukünftige Projekte nehme ich drei Dinge mit: Integrationen mit externen Systemen immer mit einem Spike starten. Automatisierte Tests von Anfang an schreiben. Und auch als Solo-Entwickler Code-Reviews simulieren – den eigenen Code nach einem Tag nochmal lesen.""",

    24: """Meine Empfehlung: Sprintify produktiv einsetzen in der Abteilung Public Cloud Solutions.

Die nächsten Schritte wären: Erstens, die Jira-Anbindung über einen schreibgeschützten Service-Account absichern. Zweitens, eine Einführungsschulung für das Team. Drittens, Feedback sammeln in den ersten drei Sprints. Viertens, automatisierte Tests ausbauen. Und mittelfristig eine mobile Ansicht für Stand-ups.

Ohne Sprintify bleibt die Sprint-Planung auf Excel angewiesen, Schätzungen basieren auf Bauchgefühl, und Sprint-Ziele werden regelmässig verfehlt.""",

    25: """Zusammengefasst: Sprintify gibt dem Team eine datenbasierte Antwort auf die Frage: Wie viele Story Points schaffen wir im nächsten Sprint?

Vielen Dank für Ihre Aufmerksamkeit. Ich bin offen für Fragen.""",
}

for idx, note_text in notes.items():
    slide = prs.slides[idx]
    if not slide.has_notes_slide:
        slide.notes_slide  # creates notes slide
    slide.notes_slide.notes_text_frame.text = note_text

prs.save("/Users/timohaldi/Documents/Github/Private/timohone/sprintify/docs/Sprintify_Praesentation.pptx")
print("Speaker notes added to all 26 slides.")
